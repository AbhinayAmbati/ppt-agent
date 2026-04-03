#!/usr/bin/env python3
"""
PPT MCP Server  (1.ppt_server.py)
==================================
MCP server that creates and styles PowerPoint files using python-pptx.
Each presentation uses a randomly selected design style from 6 distinct templates.
Content slides rotate through different layouts so no two slides look the same.

Design Styles (picked randomly per presentation)
-------------------------------------------------
1. modern_bar    — colored top title bar, full-width bullets, bottom strip
2. split_panel   — left colored panel (30%) with title, right white area with bullets
3. gradient_card — large colored number card top-left, title beside it, bullets below
4. minimal_line  — white bg, thin accent underline below title, clean bullets
5. bold_header   — full-width dark header with large title, light bg for bullets
6. sidebar       — right sidebar accent strip, bullets left with oversized slide number

Content Layout Rotation (per slide within a presentation)
----------------------------------------------------------
Slides rotate through 4 sub-layouts:
  A. full_bullets       — standard full-width bullet list
  B. two_col_bullets    — two columns of bullets side by side
  C. bullets_with_img   — bullets left + image placeholder right
  D. highlight_box      — key point in large accent box at top, rest as bullets

Themes: ocean | corporate | minimal | dark | academic
Transport: Binary stdio, newline-delimited JSON
"""

import sys
import json
import asyncio
import random
import traceback as _tb
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PROJECT_ROOT = Path(__file__).parent.parent.parent
OUTPUTS_DIR  = PROJECT_ROOT / "outputs"

_current_presentation = None
_active_theme         = None
_active_style         = None   # chosen once per presentation, used for all slides

BULLET_SYMBOLS = ["▸", "◆", "●", "◉", "▹"]

# ── THEMES ────────────────────────────────────────────────────────────────────
THEMES = {
    "ocean": {
        "bg":             RGBColor(0x0A, 0x29, 0x4A),
        "title_slide_bg": RGBColor(0x05, 0x14, 0x2E),
        "panel":          RGBColor(0x00, 0x4E, 0x7C),
        "accent":         RGBColor(0x00, 0xB4, 0xD8),
        "accent2":        RGBColor(0x90, 0xE0, 0xEF),
        "title_fg":       RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_fg":    RGBColor(0x90, 0xE0, 0xEF),
        "bullet_fg":      RGBColor(0xE0, 0xF7, 0xFF),
        "dark_text":      RGBColor(0x0A, 0x29, 0x4A),
        "img_bg":         RGBColor(0x05, 0x3A, 0x5E),
        "title_font":     "Calibri",
        "body_font":      "Calibri",
    },
    "corporate": {
        "bg":             RGBColor(0x1A, 0x1A, 0x2E),
        "title_slide_bg": RGBColor(0x10, 0x10, 0x20),
        "panel":          RGBColor(0x2A, 0x0A, 0x0A),
        "accent":         RGBColor(0xE9, 0x4F, 0x37),
        "accent2":        RGBColor(0xF5, 0xA6, 0x23),
        "title_fg":       RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_fg":    RGBColor(0xF5, 0xA6, 0x23),
        "bullet_fg":      RGBColor(0xF0, 0xF0, 0xF0),
        "dark_text":      RGBColor(0x1A, 0x1A, 0x2E),
        "img_bg":         RGBColor(0x2E, 0x1A, 0x1A),
        "title_font":     "Calibri",
        "body_font":      "Calibri",
    },
    "minimal": {
        "bg":             RGBColor(0xF8, 0xF9, 0xFA),
        "title_slide_bg": RGBColor(0xFF, 0xFF, 0xFF),
        "panel":          RGBColor(0xE3, 0xF2, 0xFD),
        "accent":         RGBColor(0x21, 0x96, 0xF3),
        "accent2":        RGBColor(0x64, 0xB5, 0xF6),
        "title_fg":       RGBColor(0x1A, 0x1A, 0x2E),
        "subtitle_fg":    RGBColor(0x21, 0x96, 0xF3),
        "bullet_fg":      RGBColor(0x2C, 0x2C, 0x3E),
        "dark_text":      RGBColor(0x1A, 0x1A, 0x2E),
        "img_bg":         RGBColor(0xDD, 0xEE, 0xFF),
        "title_font":     "Calibri",
        "body_font":      "Calibri",
    },
    "dark": {
        "bg":             RGBColor(0x12, 0x12, 0x12),
        "title_slide_bg": RGBColor(0x08, 0x08, 0x08),
        "panel":          RGBColor(0x1E, 0x1E, 0x2E),
        "accent":         RGBColor(0xBB, 0x86, 0xFC),
        "accent2":        RGBColor(0x03, 0xDA, 0xC6),
        "title_fg":       RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_fg":    RGBColor(0xBB, 0x86, 0xFC),
        "bullet_fg":      RGBColor(0xE0, 0xE0, 0xE0),
        "dark_text":      RGBColor(0x12, 0x12, 0x12),
        "img_bg":         RGBColor(0x1E, 0x1E, 0x2E),
        "title_font":     "Calibri",
        "body_font":      "Calibri",
    },
    "academic": {
        "bg":             RGBColor(0x1B, 0x2A, 0x4A),
        "title_slide_bg": RGBColor(0x0D, 0x1B, 0x33),
        "panel":          RGBColor(0x0D, 0x1B, 0x33),
        "accent":         RGBColor(0xC8, 0xA2, 0x00),
        "accent2":        RGBColor(0xE8, 0xD0, 0x70),
        "title_fg":       RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_fg":    RGBColor(0xC8, 0xA2, 0x00),
        "bullet_fg":      RGBColor(0xF0, 0xEC, 0xD8),
        "dark_text":      RGBColor(0x1B, 0x2A, 0x4A),
        "img_bg":         RGBColor(0x0D, 0x1B, 0x33),
        "title_font":     "Georgia",
        "body_font":      "Calibri",
    },
}

DEFAULT_THEME = "ocean"

# 6 distinct design styles — one is chosen randomly per presentation
DESIGN_STYLES = [
    "modern_bar",
    "split_panel",
    "gradient_card",
    "minimal_line",
    "bold_header",
    "sidebar",
]

# Content sub-layouts rotate per slide
CONTENT_LAYOUTS = ["full_bullets", "two_col_bullets", "bullets_with_img", "highlight_box"]


# ── HELPERS ───────────────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, x, y, w, h, color: RGBColor, border_color=None, border_pt=1):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, x, y, w, h, text, color, size_pt,
                 bold=False, font="Calibri", align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, word_wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.color.rgb = color
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = font
    return tb, tf


def _styled_run(para, text, color, size_pt, bold=False, font="Calibri"):
    run = para.add_run()
    run.text = str(text)
    run.font.color.rgb = color
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = font
    return run


def _remove_placeholders(slide):
    sp_tree = slide.shapes._spTree
    for ph in list(slide.placeholders):
        sp_tree.remove(ph._element)


def _bullets_tf(slide, x, y, w, h, bullets, theme, sym_idx=0):
    """Render bullet list into a text box — shared by all layouts."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    sym = BULLET_SYMBOLS[sym_idx % len(BULLET_SYMBOLS)]
    for i, text in enumerate(bullets):
        text = str(text).strip() if text else ""
        if not text:
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(9)
        p.space_after  = Pt(3)
        _styled_run(p, f"{sym}  ", theme["accent"], 15, bold=True, font=theme["body_font"])
        _styled_run(p, text, theme["bullet_fg"], 17, font=theme["body_font"])
    return tf


def _img_placeholder(slide, theme, x, y, w, h, caption=""):
    """Styled image placeholder box — used across all layouts."""
    _add_rect(slide, x, y, w, h, theme["img_bg"], border_color=theme["accent2"])
    # icon
    tb, tf = _add_textbox(slide, x+w/2-0.6, y+h/2-0.8, 1.2, 0.8,
                           "🖼", theme["accent2"], 32,
                           font="Segoe UI Emoji", align=PP_ALIGN.CENTER)
    # caption
    tb2, tf2 = _add_textbox(slide, x+0.2, y+h/2+0.05, w-0.4, 0.7,
                             f"[{caption}]" if caption else "[Image]",
                             theme["accent2"], 12, align=PP_ALIGN.CENTER)
    # hint
    tb3, tf3 = _add_textbox(slide, x+0.2, y+h-0.45, w-0.4, 0.35,
                             "Replace with actual image",
                             theme["bullet_fg"], 8, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# TITLE SLIDE STYLES  (6 distinct looks)
# ══════════════════════════════════════════════════════════════════════════════

def _title_modern_bar(slide, title, subtitle, theme):
    """Left accent bar | bottom strip | corner square | large left-aligned title."""
    _set_bg(slide, theme["title_slide_bg"])
    _remove_placeholders(slide)
    _add_rect(slide, 0, 0, 0.18, 7.5, theme["accent"])
    _add_rect(slide, 0, 6.6, 13.33, 0.9, theme["accent"])
    _add_rect(slide, 12.4, 0, 0.93, 0.9, theme["accent2"])
    _add_textbox(slide, 0.5, 1.5, 12.0, 2.5, title,
                 theme["title_fg"], 46, bold=True, font=theme["title_font"])
    _add_rect(slide, 0.5, 4.2, 4.5, 0.06, theme["accent2"])
    _add_textbox(slide, 0.5, 4.4, 11.5, 1.4, subtitle,
                 theme["subtitle_fg"], 22, font=theme["body_font"])
    _add_textbox(slide, 0.4, 6.65, 9.0, 0.5, "AUTO-GENERATED PRESENTATION",
                 theme["title_fg"], 9, bold=True, font=theme["body_font"])


def _title_split_panel(slide, title, subtitle, theme):
    """Left dark panel (40%) with title, right lighter area with subtitle."""
    _set_bg(slide, theme["bg"])
    _remove_placeholders(slide)
    _add_rect(slide, 0, 0, 5.5, 7.5, theme["panel"])
    _add_rect(slide, 5.5, 0, 7.83, 7.5, theme["title_slide_bg"])
    # vertical accent line at split
    _add_rect(slide, 5.35, 0, 0.15, 7.5, theme["accent"])
    # title on left panel
    _add_textbox(slide, 0.3, 1.0, 4.9, 4.0, title,
                 theme["title_fg"], 36, bold=True, font=theme["title_font"],
                 anchor=MSO_ANCHOR.MIDDLE)
    # accent dot
    _add_rect(slide, 0.3, 6.2, 0.6, 0.6, theme["accent"])
    # subtitle on right
    _add_textbox(slide, 5.8, 2.5, 7.0, 2.5, subtitle,
                 theme["subtitle_fg"], 24, font=theme["body_font"],
                 anchor=MSO_ANCHOR.MIDDLE)
    _add_textbox(slide, 5.8, 6.8, 6.0, 0.4, "AUTO-GENERATED PRESENTATION",
                 theme["accent2"], 9, bold=True, font=theme["body_font"])


def _title_gradient_card(slide, title, subtitle, theme):
    """Large accent block top-half, subtitle on dark bottom-half."""
    _set_bg(slide, theme["title_slide_bg"])
    _remove_placeholders(slide)
    _add_rect(slide, 0, 0, 13.33, 4.2, theme["accent"])
    # diagonal decorative strip
    _add_rect(slide, 0, 3.9, 13.33, 0.5, theme["accent2"])
    _add_textbox(slide, 0.6, 0.5, 12.0, 3.2, title,
                 theme["title_slide_bg"], 44, bold=True, font=theme["title_font"],
                 anchor=MSO_ANCHOR.MIDDLE)
    _add_textbox(slide, 0.6, 4.6, 12.0, 1.5, subtitle,
                 theme["subtitle_fg"], 24, font=theme["body_font"])
    _add_rect(slide, 0.6, 6.5, 2.0, 0.06, theme["accent2"])
    _add_textbox(slide, 0.6, 6.65, 9.0, 0.5, "AUTO-GENERATED PRESENTATION",
                 theme["accent2"], 9, bold=True, font=theme["body_font"])


def _title_minimal_line(slide, title, subtitle, theme):
    """Clean white/dark bg, accent underline below title, very minimal."""
    _set_bg(slide, theme["title_slide_bg"])
    _remove_placeholders(slide)
    # top thin accent bar
    _add_rect(slide, 0, 0, 13.33, 0.12, theme["accent"])
    _add_textbox(slide, 1.0, 1.8, 11.3, 2.4, title,
                 theme["title_fg"], 46, bold=True, font=theme["title_font"])
    # underline
    _add_rect(slide, 1.0, 4.3, 8.0, 0.08, theme["accent"])
    _add_rect(slide, 1.0, 4.5, 2.5, 0.08, theme["accent2"])
    _add_textbox(slide, 1.0, 4.7, 11.3, 1.5, subtitle,
                 theme["subtitle_fg"], 22, font=theme["body_font"])
    # bottom thin bar
    _add_rect(slide, 0, 7.38, 13.33, 0.12, theme["accent2"])
    _add_textbox(slide, 1.0, 7.05, 8.0, 0.35, "AUTO-GENERATED PRESENTATION",
                 theme["accent2"], 9, bold=True, font=theme["body_font"])


def _title_bold_header(slide, title, subtitle, theme):
    """Full-width bold accent header, subtitle on clean lower area."""
    _set_bg(slide, theme["title_slide_bg"])
    _remove_placeholders(slide)
    _add_rect(slide, 0, 0, 13.33, 3.8, theme["accent"])
    _add_textbox(slide, 0.5, 0.3, 12.3, 3.2, title,
                 theme["title_slide_bg"], 42, bold=True, font=theme["title_font"],
                 anchor=MSO_ANCHOR.MIDDLE)
    _add_rect(slide, 0, 3.8, 13.33, 0.12, theme["accent2"])
    _add_textbox(slide, 0.5, 4.1, 12.0, 1.8, subtitle,
                 theme["subtitle_fg"], 24, font=theme["body_font"])
    _add_rect(slide, 0, 6.9, 13.33, 0.6, theme["panel"])
    _add_textbox(slide, 0.5, 6.95, 9.0, 0.45, "AUTO-GENERATED PRESENTATION",
                 theme["accent2"], 9, bold=True, font=theme["body_font"])


def _title_sidebar(slide, title, subtitle, theme):
    """Right accent sidebar strip, title left-centred, subtitle below."""
    _set_bg(slide, theme["title_slide_bg"])
    _remove_placeholders(slide)
    _add_rect(slide, 12.73, 0, 0.6, 7.5, theme["accent"])
    _add_rect(slide, 12.33, 0, 0.4, 7.5, theme["accent2"])
    _add_rect(slide, 0, 6.8, 12.33, 0.7, theme["panel"])
    _add_textbox(slide, 0.6, 1.3, 11.5, 2.8, title,
                 theme["title_fg"], 44, bold=True, font=theme["title_font"])
    _add_rect(slide, 0.6, 4.2, 6.0, 0.07, theme["accent"])
    _add_textbox(slide, 0.6, 4.4, 11.0, 1.5, subtitle,
                 theme["subtitle_fg"], 22, font=theme["body_font"])
    _add_textbox(slide, 0.6, 6.85, 8.0, 0.45, "AUTO-GENERATED PRESENTATION",
                 theme["accent2"], 9, bold=True, font=theme["body_font"])


# ══════════════════════════════════════════════════════════════════════════════
# CONTENT SLIDE STYLES  (each design style × 4 content sub-layouts)
# ══════════════════════════════════════════════════════════════════════════════

def _content_modern_bar(slide, title, bullets, theme, slide_num, content_layout):
    _set_bg(slide, theme["bg"])
    _remove_placeholders(slide)
    _add_rect(slide, 0, 0, 13.33, 1.45, theme["accent"])
    strip = theme["accent2"] if slide_num % 2 == 0 else theme["accent"]
    _add_rect(slide, 0, 7.1, 13.33, 0.4, strip)
    # slide num pill
    pill = _add_rect(slide, 12.5, 0.45, 0.6, 0.45, theme["bg"])
    p = pill.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _styled_run(p, str(slide_num), theme["accent2"], 11, bold=True, font=theme["title_font"])
    _add_textbox(slide, 0.35, 0.08, 12.0, 1.28, title,
                 theme["title_fg"], 30, bold=True, font=theme["title_font"],
                 anchor=MSO_ANCHOR.MIDDLE)
    _render_content_layout(slide, bullets, theme, content_layout, top=1.6, h=5.2)


def _content_split_panel(slide, title, bullets, theme, slide_num, content_layout):
    _set_bg(slide, theme["bg"])
    _remove_placeholders(slide)
    # left panel
    _add_rect(slide, 0, 0, 3.8, 7.5, theme["panel"])
    _add_rect(slide, 3.8, 0, 0.12, 7.5, theme["accent"])
    # slide number large on panel
    _add_textbox(slide, 0.2, 0.2, 3.3, 1.0, f"{slide_num:02d}",
                 theme["accent"], 52, bold=True, font=theme["title_font"],
                 anchor=MSO_ANCHOR.TOP)
    # title on panel
    _add_textbox(slide, 0.2, 1.3, 3.3, 5.5, title,
                 theme["title_fg"], 22, bold=True, font=theme["title_font"],
                 anchor=MSO_ANCHOR.MIDDLE)
    # bottom accent dot on panel
    _add_rect(slide, 0.3, 6.9, 0.5, 0.5, theme["accent2"])
    # bullets on right
    _render_content_layout(slide, bullets, theme, content_layout,
                            top=0.6, h=6.5, left=4.2, w=8.8)


def _content_gradient_card(slide, title, bullets, theme, slide_num, content_layout):
    _set_bg(slide, theme["title_slide_bg"])
    _remove_placeholders(slide)
    # full-width accent top header
    _add_rect(slide, 0, 0, 13.33, 2.0, theme["accent"])
    _add_textbox(slide, 0.4, 0.1, 12.0, 1.7, title,
                 theme["title_slide_bg"], 30, bold=True, font=theme["title_font"],
                 anchor=MSO_ANCHOR.MIDDLE)
    # accent2 divider
    _add_rect(slide, 0, 2.0, 13.33, 0.12, theme["accent2"])
    # slide num badge top-right
    badge = _add_rect(slide, 12.1, 0.55, 0.9, 0.9, theme["title_slide_bg"])
    bp = badge.text_frame.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    _styled_run(bp, str(slide_num), theme["accent2"], 18, bold=True, font=theme["title_font"])
    # set bullet_fg to light for dark bg
    _render_content_layout(slide, bullets, theme, content_layout, top=2.3, h=4.8)


def _content_minimal_line(slide, title, bullets, theme, slide_num, content_layout):
    # white / light background with minimal chrome
    light_bg = RGBColor(0xF5, 0xF5, 0xF5) if theme["bg"] == RGBColor(0xF8, 0xF9, 0xFA) else theme["bg"]
    _set_bg(slide, light_bg)
    _remove_placeholders(slide)
    _add_rect(slide, 0, 0, 13.33, 0.08, theme["accent"])
    _add_textbox(slide, 0.6, 0.2, 11.5, 1.2, title,
                 theme["title_fg"], 32, bold=True, font=theme["title_font"])
    # double underline
    _add_rect(slide, 0.6, 1.45, 9.0, 0.07, theme["accent"])
    _add_rect(slide, 0.6, 1.6, 3.5, 0.04, theme["accent2"])
    # slide num right
    _add_textbox(slide, 12.3, 0.3, 0.8, 0.8, str(slide_num),
                 theme["accent"], 22, bold=True, font=theme["title_font"],
                 align=PP_ALIGN.CENTER)
    _add_rect(slide, 0, 7.4, 13.33, 0.1, theme["accent2"])
    _render_content_layout(slide, bullets, theme, content_layout, top=1.8, h=5.3)


def _content_bold_header(slide, title, bullets, theme, slide_num, content_layout):
    _set_bg(slide, theme["bg"])
    _remove_placeholders(slide)
    _add_rect(slide, 0, 0, 13.33, 2.1, theme["accent2"])
    _add_rect(slide, 0, 0, 0.5, 2.1, theme["accent"])   # left accent on header
    _add_textbox(slide, 0.7, 0.1, 11.8, 1.9, title,
                 theme["dark_text"], 30, bold=True, font=theme["title_font"],
                 anchor=MSO_ANCHOR.MIDDLE)
    # slide number circle top-right
    circ = _add_rect(slide, 12.4, 2.2, 0.7, 0.7, theme["accent"])
    cp = circ.text_frame.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    _styled_run(cp, str(slide_num), theme["title_fg"], 13, bold=True, font=theme["title_font"])
    _add_rect(slide, 0, 7.15, 13.33, 0.35, theme["accent"])
    _render_content_layout(slide, bullets, theme, content_layout, top=2.3, h=4.6)


def _content_sidebar(slide, title, bullets, theme, slide_num, content_layout):
    _set_bg(slide, theme["bg"])
    _remove_placeholders(slide)
    # right sidebar
    _add_rect(slide, 12.73, 0, 0.6, 7.5, theme["accent"])
    _add_rect(slide, 12.33, 0, 0.4, 7.5, theme["panel"])
    # top title bar (slightly shorter than full width)
    _add_rect(slide, 0, 0, 12.33, 1.45, theme["accent"])
    _add_textbox(slide, 0.35, 0.08, 11.7, 1.28, title,
                 theme["title_fg"], 30, bold=True, font=theme["title_font"],
                 anchor=MSO_ANCHOR.MIDDLE)
    # large slide num on sidebar
    _add_textbox(slide, 12.35, 3.2, 0.55, 1.0, str(slide_num),
                 theme["accent2"], 18, bold=True, font=theme["title_font"],
                 align=PP_ALIGN.CENTER)
    _add_rect(slide, 0, 7.15, 12.33, 0.35, theme["accent2"])
    _render_content_layout(slide, bullets, theme, content_layout,
                            top=1.6, h=5.2, left=0.4, w=11.7)


# ══════════════════════════════════════════════════════════════════════════════
# CONTENT SUB-LAYOUT RENDERER
# ══════════════════════════════════════════════════════════════════════════════

def _render_content_layout(slide, bullets, theme, layout,
                            top=1.6, h=5.2, left=0.4, w=12.5):
    """
    Render bullet content in one of 4 sub-layouts.
    All dimensions are in inches relative to passed top/h/left/w bounds.
    """
    if layout == "full_bullets":
        _bullets_tf(slide, left, top, w, h, bullets, theme)

    elif layout == "two_col_bullets":
        mid = len(bullets) // 2 + len(bullets) % 2
        col_w = (w - 0.3) / 2
        _bullets_tf(slide, left, top, col_w, h, bullets[:mid], theme, sym_idx=0)
        _bullets_tf(slide, left + col_w + 0.3, top, col_w, h, bullets[mid:], theme, sym_idx=2)

    elif layout == "bullets_with_img":
        bul_w = w * 0.57
        img_x = left + bul_w + 0.2
        img_w = w - bul_w - 0.2
        _bullets_tf(slide, left, top, bul_w, h, bullets, theme)
        _img_placeholder(slide, theme, img_x, top, img_w, h, caption="Visual")

    elif layout == "highlight_box":
        # first bullet in large accent highlight box, rest as normal bullets
        if bullets:
            box_h = 1.4
            _add_rect(slide, left, top, w, box_h, theme["accent"])
            _add_textbox(slide, left + 0.2, top + 0.1, w - 0.4, box_h - 0.15,
                         bullets[0], theme["title_slide_bg"], 18, bold=True,
                         font=theme["body_font"], anchor=MSO_ANCHOR.MIDDLE)
            if len(bullets) > 1:
                _bullets_tf(slide, left, top + box_h + 0.2, w,
                             h - box_h - 0.2, bullets[1:], theme)


# ══════════════════════════════════════════════════════════════════════════════
# CONCLUSION SLIDE STYLES
# ══════════════════════════════════════════════════════════════════════════════

def _conclusion_modern_bar(slide, title, bullets, theme, slide_num):
    _set_bg(slide, theme["title_slide_bg"])
    _remove_placeholders(slide)
    _add_rect(slide, 0, 0, 13.33, 0.15, theme["accent"])
    _add_rect(slide, 0, 7.35, 13.33, 0.15, theme["accent"])
    _add_textbox(slide, 0.5, 0.3, 12.3, 1.5, title,
                 theme["accent"], 36, bold=True, font=theme["title_font"],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_rect(slide, 2.0, 1.9, 9.33, 0.06, theme["accent2"])
    _bullets_tf(slide, 1.5, 2.1, 10.33, 4.2, bullets, theme)
    pill = _add_rect(slide, 12.55, 0.2, 0.55, 0.45, theme["accent"])
    p = pill.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _styled_run(p, str(slide_num), theme["title_fg"], 11, bold=True, font=theme["title_font"])


def _conclusion_split_panel(slide, title, bullets, theme, slide_num):
    _set_bg(slide, theme["bg"])
    _remove_placeholders(slide)
    _add_rect(slide, 0, 0, 13.33, 7.5, theme["panel"])
    _add_rect(slide, 1.5, 1.0, 10.33, 0.1, theme["accent"])
    _add_textbox(slide, 1.5, 1.2, 10.33, 1.5, title,
                 theme["accent"], 38, bold=True, font=theme["title_font"],
                 align=PP_ALIGN.CENTER)
    _add_rect(slide, 1.5, 2.8, 10.33, 0.06, theme["accent2"])
    _bullets_tf(slide, 2.0, 3.0, 9.33, 3.8, bullets, theme)
    _add_textbox(slide, 0.5, 6.9, 5.0, 0.4, f"Slide {slide_num}",
                 theme["accent2"], 10, font=theme["body_font"])


def _conclusion_gradient_card(slide, title, bullets, theme, slide_num):
    _set_bg(slide, theme["title_slide_bg"])
    _remove_placeholders(slide)
    _add_rect(slide, 0, 0, 13.33, 2.8, theme["accent2"])
    _add_textbox(slide, 0.5, 0.2, 12.3, 2.4, title,
                 theme["dark_text"], 36, bold=True, font=theme["title_font"],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_rect(slide, 0, 2.8, 13.33, 0.15, theme["accent"])
    _bullets_tf(slide, 1.5, 3.2, 10.33, 3.8, bullets, theme)


# Style dispatch tables
_TITLE_BUILDERS = {
    "modern_bar":    _title_modern_bar,
    "split_panel":   _title_split_panel,
    "gradient_card": _title_gradient_card,
    "minimal_line":  _title_minimal_line,
    "bold_header":   _title_bold_header,
    "sidebar":       _title_sidebar,
}

_CONTENT_BUILDERS = {
    "modern_bar":    _content_modern_bar,
    "split_panel":   _content_split_panel,
    "gradient_card": _content_gradient_card,
    "minimal_line":  _content_minimal_line,
    "bold_header":   _content_bold_header,
    "sidebar":       _content_sidebar,
}

_CONCLUSION_BUILDERS = {
    "modern_bar":    _conclusion_modern_bar,
    "split_panel":   _conclusion_split_panel,
    "gradient_card": _conclusion_gradient_card,
    "minimal_line":  _conclusion_modern_bar,   # reuse
    "bold_header":   _conclusion_split_panel,  # reuse
    "sidebar":       _conclusion_gradient_card,# reuse
}


# ── TOOL FUNCTIONS ────────────────────────────────────────────────────────────

async def create_presentation(title: str, subtitle: str = "",
                               theme_name: str = DEFAULT_THEME) -> dict:
    """
    Initialize a new 16:9 Presentation, pick a random design style,
    and render the title slide.
    """
    global _current_presentation, _active_theme, _active_style
    theme = THEMES.get(theme_name, THEMES[DEFAULT_THEME])
    _active_theme = theme
    _active_style = random.choice(DESIGN_STYLES)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _TITLE_BUILDERS[_active_style](slide, title, subtitle, theme)

    _current_presentation = prs
    return {
        "status": "success", "message": "Presentation created",
        "slide_count": 1, "design_style": _active_style,
    }


async def add_slide(layout_type: str = "title_and_content") -> dict:
    """Append a blank slide. Returns slide_count so caller can derive the new index."""
    global _current_presentation
    if _current_presentation is None:
        return {"status": "error", "message": "No active presentation"}
    blank = _current_presentation.slide_layouts[6]
    _current_presentation.slides.add_slide(blank)
    return {"status": "success", "message": "Slide added",
            "slide_count": len(_current_presentation.slides)}


async def write_text_to_slide(slide_index: int, title: str, content: list,
                               include_image: bool = False,
                               image_path: str = None,
                               is_conclusion: bool = False) -> dict:
    """
    Style and populate a slide using the presentation's active design style.
    Content sub-layout rotates: full_bullets → two_col_bullets → bullets_with_img
    → highlight_box → repeat.
    """
    global _current_presentation, _active_theme, _active_style
    if _current_presentation is None:
        return {"status": "error", "message": "No active presentation"}

    slides = _current_presentation.slides
    if slide_index < 0 or slide_index >= len(slides):
        return {"status": "error",
                "message": f"slide_index {slide_index} out of range (total={len(slides)})"}

    clean = [str(b).strip() for b in (content or []) if b is not None and str(b).strip()]
    if not clean:
        clean = [f"Key information about {title}"]

    slide = slides[slide_index]
    style = _active_style or "modern_bar"
    theme = _active_theme or THEMES[DEFAULT_THEME]

    try:
        if is_conclusion:
            _CONCLUSION_BUILDERS[style](slide, str(title), clean, theme, slide_index)
        else:
            # rotate content sub-layout based on slide index
            layouts = CONTENT_LAYOUTS
            if include_image:
                # force image layout for slides marked for placeholder
                sub = "bullets_with_img"
            else:
                sub = layouts[(slide_index - 1) % len(layouts)]
            _CONTENT_BUILDERS[style](slide, str(title), clean, theme, slide_index, sub)
    except Exception as exc:
        return {"status": "error", "message": f"Styling failed: {exc}",
                "trace": _tb.format_exc()}

    return {"status": "success", "message": "Text written", "bullet_count": len(clean),
            "design_style": style}


async def add_image_placeholder(slide_index: int, placeholder_text: str = "[Image]") -> dict:
    """Add a standalone styled image placeholder box to an existing slide."""
    global _current_presentation, _active_theme
    if _current_presentation is None:
        return {"status": "error", "message": "No active presentation"}
    slides = _current_presentation.slides
    if slide_index < 0 or slide_index >= len(slides):
        return {"status": "error", "message": f"slide_index {slide_index} out of range"}
    _img_placeholder(slides[slide_index], _active_theme or THEMES[DEFAULT_THEME],
                     x=1.5, y=2.0, w=10.33, h=4.0, caption=placeholder_text)
    return {"status": "success", "message": "Image placeholder added"}


async def set_theme(theme_name: str) -> dict:
    """Switch the active theme for subsequent slides."""
    global _active_theme
    if theme_name not in THEMES:
        return {"status": "error", "message": f"Unknown theme. Available: {list(THEMES.keys())}"}
    _active_theme = THEMES[theme_name]
    return {"status": "success", "theme": theme_name}


async def save_presentation(file_path: str) -> dict:
    """Save the in-memory Presentation to outputs/."""
    global _current_presentation
    if _current_presentation is None:
        return {"status": "error", "message": "No active presentation"}
    OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)
    full_path = OUTPUTS_DIR / Path(file_path).name
    _current_presentation.save(str(full_path))
    return {"status": "success", "message": "Saved", "file_path": str(full_path)}


async def get_presentation_info() -> dict:
    """Return metadata about the current in-memory presentation."""
    if _current_presentation is None:
        return {"status": "error", "message": "No active presentation"}
    return {"status": "success", "slide_count": len(_current_presentation.slides),
            "design_style": _active_style}


async def call_tool(name: str, arguments: dict) -> dict:
    """Route an MCP tool-call to the appropriate handler."""
    try:
        if name == "create_presentation":
            return await create_presentation(arguments["title"], arguments.get("subtitle", ""),
                                              arguments.get("theme_name", DEFAULT_THEME))
        elif name == "add_slide":
            return await add_slide(arguments.get("layout_type", "title_and_content"))
        elif name == "write_text_to_slide":
            return await write_text_to_slide(
                arguments["slide_index"], arguments["title"], arguments["content"],
                include_image=arguments.get("include_image", False),
                image_path=arguments.get("image_path", None),
                is_conclusion=arguments.get("is_conclusion", False))
        elif name == "add_image_placeholder":
            return await add_image_placeholder(arguments["slide_index"],
                                                arguments.get("placeholder_text", "[Image]"))
        elif name == "save_presentation":
            return await save_presentation(arguments["file_path"])
        elif name == "get_presentation_info":
            return await get_presentation_info()
        elif name == "set_theme":
            return await set_theme(arguments["theme_name"])
        else:
            return {"status": "error", "message": f"Unknown tool: {name}"}
    except Exception as e:
        return {"status": "error", "message": str(e), "trace": _tb.format_exc()}


async def main():
    """Binary stdio event loop."""
    stdin  = sys.stdin.buffer
    stdout = sys.stdout.buffer
    while True:
        line = stdin.readline()
        if not line:
            break
        try:
            req = json.loads(line.decode("utf-8").strip())
            if req.get("method") == "tools/call":
                result = await call_tool(req["params"]["name"], req["params"].get("arguments", {}))
            else:
                result = {"status": "error", "message": f"Unknown method: {req.get('method')}"}
        except Exception as e:
            result = {"status": "error", "message": f"Parse error: {e}"}
        stdout.write((json.dumps(result) + "\n").encode("utf-8"))
        stdout.flush()


if __name__ == "__main__":
    asyncio.run(main())
